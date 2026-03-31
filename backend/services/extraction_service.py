import base64
import io

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _image_to_png_base64(image_bytes: bytes) -> tuple[str, int, int]:
    """Convert image bytes to PNG base64 string, return (base64, width, height)."""
    img = Image.open(io.BytesIO(image_bytes))
    img = img.convert("RGBA")
    w, h = img.size
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    b64 = base64.b64encode(buf.getvalue()).decode("ascii")
    return b64, w, h


def extract_images_from_pdf(file_bytes: bytes) -> list[dict]:
    """Extract all embedded images from a PDF file."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    layers = []
    seen_xrefs = set()

    for page_num in range(len(doc)):
        page = doc[page_num]
        images = page.get_images(full=True)

        for img_index, img_info in enumerate(images):
            xref = img_info[0]
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)

            extracted = doc.extract_image(xref)
            if not extracted or not extracted.get("image"):
                continue

            try:
                b64, w, h = _image_to_png_base64(extracted["image"])
            except Exception:
                continue

            # Skip very small images (icons, bullets, etc.)
            if w < 50 or h < 50:
                continue

            layers.append({
                "name": f"PDF 圖片 {len(layers) + 1}",
                "image_base64": b64,
                "bounds": {"x": 0, "y": 0, "w": w, "h": h},
            })

    doc.close()

    if not layers:
        raise ValueError("PDF 中未找到任何圖片")

    return layers


def extract_images_from_pptx(file_bytes: bytes) -> list[dict]:
    """Extract all images from a PPTX presentation."""
    prs = Presentation(io.BytesIO(file_bytes))
    layers = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        img_num = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
                try:
                    blob = shape.image.blob
                except Exception:
                    continue

                try:
                    b64, w, h = _image_to_png_base64(blob)
                except Exception:
                    continue

                if w < 50 or h < 50:
                    continue

                img_num += 1
                # Convert EMU to pixels (96 DPI)
                x = int(shape.left / 914400 * 96) if shape.left else 0
                y = int(shape.top / 914400 * 96) if shape.top else 0

                layers.append({
                    "name": f"投影片 {slide_num} - 圖片 {img_num}",
                    "image_base64": b64,
                    "bounds": {"x": x, "y": y, "w": w, "h": h},
                })

    if not layers:
        raise ValueError("PPTX 中未找到任何圖片")

    return layers
